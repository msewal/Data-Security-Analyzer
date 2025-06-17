import io
import os
import re
import urllib.parse
import logging
import time
import json
import shutil
import hashlib
import mimetypes
import docx
import PyPDF2
import pandas as pd
import chardet
import magic
import math
from django.shortcuts import render, redirect
from django.http import JsonResponse, HttpResponse
from django.conf import settings
from django.urls import reverse
from django.views.decorators.http import require_http_methods, require_POST
from django.views.decorators.csrf import csrf_exempt
from malware.models import QuarantinedFile
from django.contrib import messages
from datetime import datetime
from typing import Dict, List, Tuple
from .patterns import (
    ALL_REGEX_PATTERNS_BACKEND,
    get_all_patterns,
    validate_regex_pattern,
    get_context_lines,
    should_scan_file,
    sensitive_patterns,
    compile_patterns
)
from collections import defaultdict
from .utils.scan_tools import scan_with_grep, scan_with_ripgrep, scan_with_ag, scan_with_ack
import math
import base64
from cryptography.fernet import Fernet
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC

# Use BASE_DIR from settings
BASE_DIR = settings.BASE_DIR

# Office ve PDF işleme için gerekli importlar
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

def get_file_content(file_path):
    """Dosya içeriğini okur ve metin olarak döndürür."""
    try:
        file_type = get_file_type(file_path)
        
        # Metin dosyaları
        if file_type in ['txt', 'md', 'log', 'ini', 'conf', 'cfg', 'properties', 'html', 'htm', 'css', 'js']:
            try:
                # Önce UTF-8 ile dene
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                # UTF-8 başarısız olursa, chardet ile encoding'i tespit et
                import chardet
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                    result = chardet.detect(raw_data)
                    encoding = result['encoding']
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
        
        # PDF dosyaları
        elif file_type == 'pdf':
            try:
                import PyPDF2
                with open(file_path, 'rb') as f:
                    pdf = PyPDF2.PdfReader(f)
                    text = ''
                    for page in pdf.pages:
                        text += page.extract_text() + '\n'
                    return text
            except Exception as e:
                logger.error(f"Error reading PDF {file_path}: {str(e)}")
                return None
        
        # Word dosyaları
        elif file_type in ['docx', 'doc']:
            try:
                import docx
                doc = docx.Document(file_path)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            except Exception as e:
                logger.error(f"Error reading Word file {file_path}: {str(e)}")
                return None
        
        # Excel dosyaları
        elif file_type in ['xlsx', 'xls']:
            try:
                import pandas as pd
                df = pd.read_excel(file_path)
                return df.to_string()
            except Exception as e:
                logger.error(f"Error reading Excel file {file_path}: {str(e)}")
                return None
        
        # CSV ve TSV dosyaları
        elif file_type in ['csv', 'tsv']:
            try:
                import pandas as pd
                sep = ',' if file_type == 'csv' else '\t'
                df = pd.read_csv(file_path, sep=sep)
                return df.to_string()
            except Exception as e:
                logger.error(f"Error reading CSV/TSV file {file_path}: {str(e)}")
                return None
        
        # JSON dosyaları
        elif file_type == 'json':
            try:
                import json
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    return json.dumps(data, indent=2, ensure_ascii=False)
            except Exception as e:
                logger.error(f"Error reading JSON file {file_path}: {str(e)}")
                return None
        
        # XML dosyaları
        elif file_type == 'xml':
            try:
                import xml.etree.ElementTree as ET
                tree = ET.parse(file_path)
                root = tree.getroot()
                return ET.tostring(root, encoding='unicode')
            except Exception as e:
                logger.error(f"Error reading XML file {file_path}: {str(e)}")
                return None
        
        # YAML dosyaları
        elif file_type in ['yaml', 'yml']:
            try:
                import yaml
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = yaml.safe_load(f)
                    return yaml.dump(data, allow_unicode=True)
            except Exception as e:
                logger.error(f"Error reading YAML file {file_path}: {str(e)}")
                return None
        
        # Arşiv dosyaları
        elif file_type in ['zip', 'rar']:
            try:
                import tempfile
                import os
                import zipfile
                import rarfile
                
                # Geçici dizin oluştur
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Arşivi aç
                    if file_type == 'zip':
                        with zipfile.ZipFile(file_path, 'r') as zip_ref:
                            zip_ref.extractall(temp_dir)
                    else:  # rar
                        with rarfile.RarFile(file_path, 'r') as rar_ref:
                            rar_ref.extractall(temp_dir)
                    
                    # Tüm dosyaları tara
                    all_content = []
                    for root, _, files in os.walk(temp_dir):
                        for file in files:
                            file_path = os.path.join(root, file)
                            content = get_file_content(file_path)
                            if content:
                                all_content.append(f"=== {file} ===\n{content}\n")
                    
                    return '\n'.join(all_content)
            except Exception as e:
                logger.error(f"Error processing archive {file_path}: {str(e)}")
                return None
        
        # Görsel dosyaları
        elif file_type in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
            try:
                import pytesseract
                from PIL import Image
                
                # Update Tesseract path for Linux
                pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'
                
                # Görüntüyü aç ve OCR uygula
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image, lang='tur+eng')
                return text
            except Exception as e:
                logger.error(f"Error processing image {file_path}: {str(e)}")
                return None
        
        else:
            logger.warning(f"Unsupported file type: {file_type}")
            return None
            
    except Exception as e:
        logger.error(f"Error reading file {file_path}: {str(e)}")
        return None

def get_file_type(file_path):
    """Dosya türünü belirler."""
    try:
        # Önce uzantıya bak
        ext = os.path.splitext(file_path)[1].lower().lstrip('.')
        
        # MIME türünü kontrol et
        try:
            import magic
            mime = magic.from_file(file_path, mime=True)
        except Exception as e:
            logger.warning(f"Error getting MIME type for {file_path}: {str(e)}")
            mime = None
        
        # Uzantı ve MIME türüne göre dosya türünü belirle
        if ext in ['txt', 'md', 'log', 'ini', 'conf', 'cfg', 'properties']:
            return 'txt'
        elif ext in ['docx', 'doc']:
            return 'docx'
        elif ext in ['xlsx', 'xls']:
            return 'xlsx'
        elif ext in ['pptx', 'ppt']:
            return 'pptx'
        elif ext == 'pdf':
            return 'pdf'
        elif ext in ['csv', 'tsv']:
            return ext
        elif ext in ['json', 'xml', 'yaml', 'yml']:
            return ext
        elif ext in ['html', 'htm', 'css', 'js']:
            return ext
        elif ext in ['zip', 'rar']:
            return ext
        elif ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
            return ext
        else:
            # MIME türüne göre kontrol et
            if mime:
                if mime.startswith('text/'):
                    return 'txt'
                elif mime == 'application/pdf':
                    return 'pdf'
                elif mime in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                            'application/msword']:
                    return 'docx'
                elif mime in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                            'application/vnd.ms-excel']:
                    return 'xlsx'
                elif mime in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                            'application/vnd.ms-powerpoint']:
                    return 'pptx'
                elif mime == 'application/json':
                    return 'json'
                elif mime == 'application/xml':
                    return 'xml'
                elif mime == 'application/zip':
                    return 'zip'
                elif mime.startswith('image/'):
                    return ext
            return None
    except Exception as e:
        logger.error(f"Error determining file type for {file_path}: {str(e)}")
        return None

def regex_search(request):
    """Regex araması yapar."""
    if request.method == 'POST':
        directory = request.POST.get('directory', '')
        file_types = request.POST.getlist('file_types', [])
        case_sensitive = request.POST.get('case_sensitive', 'false') == 'true'
        multiline = request.POST.get('multiline', 'false') == 'true'
        selected_categories = request.POST.getlist('categories', [])
        selected_subcategories = request.POST.getlist('subcategories', [])
        
        if not directory or not os.path.exists(directory):
            messages.error(request, 'Geçersiz dizin yolu.')
            return redirect('regex:regex_search')
            
        if not file_types:
            messages.error(request, 'En az bir dosya türü seçmelisiniz.')
            return redirect('regex:regex_search')
            
        if not selected_categories and not selected_subcategories:
            messages.error(request, 'En az bir kategori veya alt kategori seçmelisiniz.')
            return redirect('regex:regex_search')
            
        # Seçilen kategorilere göre regex desenlerini al
        patterns = []
        if selected_categories:
            for category in selected_categories:
                patterns.extend(ALL_REGEX_PATTERNS_BACKEND.get(category, []))
                
        if selected_subcategories:
            for subcategory in selected_subcategories:
                for category in ALL_REGEX_PATTERNS_BACKEND.items():
                    if subcategory in category[1]:
                        patterns.extend(category[1][subcategory])
                        
        if not patterns:
            messages.error(request, 'Geçerli regex deseni bulunamadı.')
            return redirect('regex:regex_search')
            
        # Regex desenlerini derle
        compiled_patterns = compile_patterns(patterns)
        
        results = {
            'matches': [],
            'errors': [],
            'skipped': []
        }
        
        for root, _, files in os.walk(directory):
            for file in files:
                file_path = os.path.join(root, file)
                file_type = get_file_type(file_path)
                
                if file_type not in file_types:
                    results['skipped'].append({
                        'path': file_path,
                        'reason': 'Desteklenmeyen dosya türü'
                    })
                    continue
                    
                content = get_file_content(file_path)
                if content is None:
                    results['errors'].append({
                        'path': file_path,
                        'error': 'Dosya okunamadı'
                    })
                    continue
                    
                file_matches = []
                for pattern in compiled_patterns:
                    matches = pattern.finditer(content)
                    for match in matches:
                        file_matches.append({
                            'pattern': pattern.pattern,
                            'match': match.group(),
                            'line': content[:match.start()].count('\n') + 1,
                            'position': match.start()
                        })
                        
                if file_matches:
                    results['matches'].append({
                        'path': file_path,
                        'matches': file_matches
                    })
                    
        # Sonuçları oturuma kaydet
        request.session['regex_results'] = results
        
        return redirect('regex:regex_search_results')
        
    # GET isteği için
    return render(request, 'regex/regex_search.html', {
        'categories': get_category_names(),
        'file_types': ALL_REGEX_PATTERNS_BACKEND.file_types
    })

def regex_search_results(request):
    """Regex arama sonuçlarını gösterir."""
    results = request.session.get('regex_results', {})
    if not results:
        messages.error(request, 'Arama sonucu bulunamadı.')
        return redirect('regex:regex_search')
        
    return render(request, 'regex/regex_search_results.html', {
        'results': results
    })

def regex_search_detail(request, file_path):
    """Belirli bir dosyanın regex arama sonuçlarını gösterir."""
    results = request.session.get('regex_results', {})
    if not results:
        messages.error(request, 'Arama sonucu bulunamadı.')
        return redirect('regex:regex_search')
        
    file_results = next(
        (result for result in results['matches'] if result['path'] == file_path),
        None
    )
    
    if not file_results:
        messages.error(request, 'Dosya sonuçları bulunamadı.')
        return redirect('regex:regex_search_results')
        
    return render(request, 'regex/regex_search_detail.html', {
        'file_path': file_path,
        'matches': file_results['matches']
    })

def quarantine_file(request, file_path):
    """Dosyayı karantinaya alır."""
    try:
        quarantine_dir = os.path.join(settings.MEDIA_ROOT, 'quarantine')
        os.makedirs(quarantine_dir, exist_ok=True)
        
        file_name = os.path.basename(file_path)
        quarantine_path = os.path.join(quarantine_dir, file_name)
        
        os.rename(file_path, quarantine_path)
        messages.success(request, f'Dosya karantinaya alındı: {file_name}')
        
    except Exception as e:
        messages.error(request, f'Karantina hatası: {str(e)}')
        
    return redirect('regex:regex_search_results')

def edit_file(request, file_path):
    """Dosya içeriğini düzenler."""
    if request.method == 'POST':
        try:
            content = request.POST.get('content', '')
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(content)
            messages.success(request, 'Dosya başarıyla güncellendi.')
            
        except Exception as e:
            messages.error(request, f'Dosya güncelleme hatası: {str(e)}')
            
        return redirect('regex:regex_search_results')
        
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            
    except Exception as e:
        messages.error(request, f'Dosya okuma hatası: {str(e)}')
        return redirect('regex:regex_search_results')
        
    return render(request, 'regex/edit_file.html', {
        'file_path': file_path,
        'content': content
    })

def is_safe_path(path):
    """Dosya yolunun güvenli olup olmadığını kontrol eder."""
    try:
        # Mutlak yol kontrolü
        abs_path = os.path.abspath(path)
        # İzin verilen dizinlerin kontrolü
        allowed_dirs = [
            os.path.abspath(settings.MEDIA_ROOT),
            os.path.abspath(settings.BASE_DIR)
        ]
        return any(abs_path.startswith(d) for d in allowed_dirs)
    except Exception:
        return False

def get_category_names():
    """Kategori isimlerini döndürür."""
    return ALL_REGEX_PATTERNS_BACKEND.get_category_names()

def sensitive_scan(request):
    """Hassas veri taraması sayfasını gösterir ve tarama sonuçlarını işler."""
    if request.method == 'POST':
        directory = request.POST.get('directory', '').strip()
        selected_categories = request.POST.getlist('categories')
        selected_subcategories = request.POST.getlist('subcategories[]')
        selected_file_types = request.POST.getlist('file_types')

        if not directory:
            return render(request, 'regex/sensitive_scan.html', {
                'error_message': 'Lütfen bir dizin seçin.',
                'categories': get_category_names(),
                'file_types': ALL_REGEX_PATTERNS_BACKEND.file_types
            })

        if not os.path.exists(directory):
            return render(request, 'regex/sensitive_scan.html', {
                'error_message': f'Dizin bulunamadı: {directory}',
                'categories': get_category_names(),
                'file_types': ALL_REGEX_PATTERNS_BACKEND.file_types
            })

        patterns = []
        if selected_categories:
            for category in selected_categories:
                patterns.extend(ALL_REGEX_PATTERNS_BACKEND.get_patterns_by_category(category))

        if selected_subcategories:
            for subcategory in selected_subcategories:
                for category in ALL_REGEX_PATTERNS_BACKEND.categories.values():
                    if subcategory in category.get('subcategories', {}):
                        patterns.extend(category['subcategories'][subcategory]['patterns'])

        if not patterns:
            return render(request, 'regex/sensitive_scan.html', {
                'error_message': 'Lütfen en az bir kategori veya alt kategori seçin.',
                'categories': get_category_names(),
                'file_types': ALL_REGEX_PATTERNS_BACKEND.file_types
            })

        results = []
        error_files = []
        skipped_files = []

        # Derlenmiş regex desenleri
        compiled_patterns = []
        for pattern in patterns:
            try:
                compiled_patterns.append(re.compile(pattern, re.IGNORECASE | re.MULTILINE))
            except re.error as e:
                logger.error(f"Regex derleme hatası: {pattern} - {str(e)}")

        for root, dirs, files in os.walk(directory):
            # Gizli dizinleri ve özel dizinleri atla
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in ['node_modules', 'venv', '__pycache__']]
            
            for file in files:
                file_path = os.path.join(root, file)
                file_ext = os.path.splitext(file)[1].lower().lstrip('.')

                # Dosya türü kontrolü
                if file_ext not in selected_file_types:
                    continue

                # Dosya boyutu kontrolü (5MB'dan büyük dosyaları atla)
                try:
                    if os.path.getsize(file_path) > 5 * 1024 * 1024:
                        skipped_files.append(f"{file_path} (Dosya boyutu çok büyük: >5MB)")
                        continue
                except OSError as e:
                    error_files.append(f"{file_path} (Dosya boyutu alınamadı: {str(e)})")
                    continue

                # Dosya içeriğini oku
                try:
                    # Önce dosyanın metin dosyası olup olmadığını kontrol et
                    with open(file_path, 'rb') as f:
                        chunk = f.read(1024)
                        if b'\x00' in chunk:  # Binary dosya kontrolü
                            skipped_files.append(f"{file_path} (Binary dosya)")
                            continue

                    # Dosyayı oku
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                        file_matches = defaultdict(lambda: defaultdict(list))

                        # Her desen için eşleşmeleri ara
                        for pattern in compiled_patterns:
                            for match in pattern.finditer(content):
                                line_number = content[:match.start()].count('\n') + 1
                                line = content.split('\n')[line_number - 1]
                                context = get_context_lines(content, line_number)

                                # Eşleşmenin hangi kategori ve alt kategoriye ait olduğunu bul
                                for category, data in ALL_REGEX_PATTERNS_BACKEND.categories.items():
                                    for subcategory, subdata in data['subcategories'].items():
                                        if pattern.pattern in subdata['patterns']:
                                            file_matches[data['name']][subdata['name']].append({
                                                'line': line_number,
                                                'match': match.group(),
                                                'context': context,
                                                'pattern': pattern.pattern
                                            })

                        if file_matches:
                            results.append({
                                'file_path': file_path,
                                'matches': file_matches
                            })

                except UnicodeDecodeError:
                    error_files.append(f"{file_path} (Karakter kodlaması hatası)")
                except PermissionError:
                    error_files.append(f"{file_path} (Erişim izni yok)")
                except Exception as e:
                    error_files.append(f"{file_path} ({str(e)})")

        return render(request, 'regex/sensitive_scan_results.html', {
            'results': results,
            'error_files': error_files,
            'skipped_files': skipped_files,
            'directory': directory,
            'stats': {
                'total_files': len(results) + len(error_files) + len(skipped_files),
                'matched_files': len(results),
                'error_files': len(error_files),
                'skipped_files': len(skipped_files)
            }
        })

    return render(request, 'regex/sensitive_scan.html', {
        'categories': get_category_names(),
        'file_types': ALL_REGEX_PATTERNS_BACKEND.file_types
    })

def sensitive_scan_results(request):
    """Hassas veri taraması sonuçlarını gösterir."""
    results = request.session.get('sensitive_scan_results', {})
    if not results:
        messages.error(request, 'Tarama sonucu bulunamadı.')
        return redirect('regex:sensitive_scan')
        
    return render(request, 'regex/sensitive_scan_results.html', {
        'results': results['results'],
        'error_files': results['error_files'],
        'skipped_files': results['skipped_files'],
        'stats': results['stats']
    })

def sensitive_scan_detail(request, file_path):
    """Hassas veri taraması sonuçlarının detaylarını gösterir."""
    # POST istekleri için maskeleme ve şifreleme işlemleri
    if request.method == 'POST':
        action = request.POST.get('action')
        password = request.POST.get('password', '')
        mask_type = request.POST.get('mask_type', 'asterisk')
        
        if action == 'mask':
            try:
                enhanced_content = get_enhanced_file_content(file_path)
                masked_content = mask_sensitive_data(enhanced_content['raw_content'], mask_type)
                return JsonResponse({
                    'success': True,
                    'masked_content': masked_content,
                    'mask_type': mask_type
                })
            except Exception as e:
                return JsonResponse({'success': False, 'error': str(e)})
                
        elif action == 'encrypt':
            try:
                if not password:
                    return JsonResponse({'success': False, 'error': 'Şifre gerekli'})
                enhanced_content = get_enhanced_file_content(file_path)
                encrypted_content = encrypt_text(enhanced_content['raw_content'], password)
                return JsonResponse({
                    'success': True,
                    'encrypted_content': encrypted_content
                })
            except Exception as e:
                return JsonResponse({'success': False, 'error': str(e)})
                
        elif action == 'decrypt':
            try:
                if not password:
                    return JsonResponse({'success': False, 'error': 'Şifre gerekli'})
                encrypted_data = request.POST.get('encrypted_data', '')
                decrypted_content = decrypt_text(encrypted_data, password)
                if decrypted_content:
                    return JsonResponse({
                        'success': True,
                        'decrypted_content': decrypted_content
                    })
                else:
                    return JsonResponse({'success': False, 'error': 'Şifre çözme başarısız'})
            except Exception as e:
                return JsonResponse({'success': False, 'error': str(e)})
    
    try:
        # Dosya bilgilerini al
        file_size = None
        file_type = None
        if os.path.exists(file_path):
            file_size = format_file_size(os.path.getsize(file_path))
            file_type = get_file_type(file_path)
        
        # Gelişmiş dosya içeriği çıkarma
        enhanced_content = get_enhanced_file_content(file_path)
        file_content = enhanced_content['raw_content']
        
        if not file_content:
            return render(request, 'regex/sensitive_scan_detail.html', {
                'file_path': file_path,
                'error_message': 'Dosya içeriği okunamadı veya boş.',
                'matches': [],
                'file_size': file_size,
                'file_type': file_type
            })

        # Eşleşmeleri bul ve içerikteki konumlarını işaretle
        matches = []
        highlighted_content = file_content
        match_positions = []  # Eşleşmelerin konumlarını sakla
        
        for category, data in ALL_REGEX_PATTERNS_BACKEND.categories.items():
            for subcategory, subdata in data['subcategories'].items():
                for pattern in subdata['patterns']:
                    try:
                        compiled_pattern = re.compile(pattern, re.IGNORECASE | re.MULTILINE)
                        for match in compiled_pattern.finditer(file_content):
                            line_number = file_content[:match.start()].count('\n') + 1
                            line = file_content.split('\n')[line_number - 1]
                            context = get_context_lines(file_content, line_number)
                            
                            matches.append({
                                'category': data['name'],
                                'subcategory': subdata['name'],
                                'line': line_number,
                                'match': match.group(),
                                'context': context,
                                'pattern': pattern,
                                'start': match.start(),
                                'end': match.end()
                            })
                            
                            # Eşleşme pozisyonunu kaydet
                            match_positions.append({
                                'start': match.start(),
                                'end': match.end(),
                                'text': match.group()
                            })
                    except re.error as e:
                        logger.error(f"Regex derleme hatası: {pattern} - {str(e)}")
                        continue

        # Eşleşmeleri satır numarasına göre sırala
        matches.sort(key=lambda x: x['line'])
        
        # Eşleşmeleri pozisyona göre sırala (tersten, böylece vurgulama işlemi doğru çalışır)
        match_positions.sort(key=lambda x: x['start'], reverse=True)
        
        # İçerikte eşleşmeleri vurgula
        for pos in match_positions:
            highlighted_text = f'<span class="highlighted-match">{pos["text"]}</span>'
            highlighted_content = (
                highlighted_content[:pos['start']] + 
                highlighted_text + 
                highlighted_content[pos['end']:]
            )
        
        # İçeriği satırlara böl
        content_lines = highlighted_content.split('\n')
        
        # İstatistikler için benzersiz değerleri hesapla
        unique_categories = list(set([m['category'] for m in matches]))
        unique_subcategories = list(set([m['subcategory'] for m in matches]))
        unique_patterns = list(set([m['pattern'] for m in matches]))

        return render(request, 'regex/sensitive_scan_detail.html', {
            'file_path': file_path,
            'matches': matches,
            'file_content': highlighted_content,
            'file_content_lines': content_lines,
            'file_size': file_size,
            'file_type': file_type,
            'file_metadata': enhanced_content.get('metadata', {}),
            'extraction_method': enhanced_content.get('extraction_method', 'basic'),
            'page_data': enhanced_content.get('pages', []),
            'unique_categories': unique_categories,
            'unique_subcategories': unique_subcategories,
            'unique_patterns': unique_patterns
        })
        
    except Exception as e:
        logger.error(f"Sensitive scan detail error: {str(e)}")
        return render(request, 'regex/sensitive_scan_detail.html', {
            'file_path': file_path,
            'error_message': f'Dosya analizi sırasında hata oluştu: {str(e)}',
            'matches': []
        })

def format_file_size(size_bytes):
    """Dosya boyutunu okunabilir formata çevirir."""
    if size_bytes == 0:
        return "0 B"
    size_names = ["B", "KB", "MB", "GB", "TB"]
    i = int(math.floor(math.log(size_bytes, 1024)))
    p = math.pow(1024, i)
    s = round(size_bytes / p, 2)
    return f"{s} {size_names[i]}"

def api_get_regex_patterns(request):
    """Regex desenlerini JSON formatında döndürür."""
    patterns = {
        'categories': ALL_REGEX_PATTERNS_BACKEND.categories,
        'file_types': ALL_REGEX_PATTERNS_BACKEND.file_types
    }
    return JsonResponse(patterns)

def quarantine_list(request):
    files = QuarantinedFile.objects.all().order_by('-quarantine_time')
    return render(request, 'regex/quarantine_list.html', {'files': files})

def view_file(request, file_path):
    """Dosyayı görüntüler ve regex eşleşmelerini vurgular."""
    try:
        # Dosya içeriğini al
        content = get_file_content(file_path)
        if content is None:
            messages.error(request, 'Dosya okunamadı.')
            return redirect('regex:regex_search_results')

        # Eşleşen regex desenlerini al
        results = request.session.get('regex_results', {})
        file_results = next(
            (result for result in results['matches'] if result['path'] == file_path),
            None
        )

        if not file_results:
            messages.error(request, 'Dosya sonuçları bulunamadı.')
            return redirect('regex:regex_search_results')

        # Eşleşmeleri vurgula
        for match in file_results['matches']:
            for category, subcategories in match.items():
                for subcategory, matches in subcategories.items():
                    for m in matches:
                        # Eşleşmeyi üstü çizili yap
                        content = content.replace(m['match'], f'<del>{m['match']}</del>')

        return render(request, 'regex/view_file.html', {
            'file_path': file_path,
            'content': content
        })

    except Exception as e:
        messages.error(request, f'Dosya görüntüleme hatası: {str(e)}')
        return redirect('regex:regex_search_results')

@csrf_exempt
def do_search(request):
    """Formdan gelen parametrelere göre tarama yapar ve sonuçları kaydeder."""
    if request.method != 'POST':
        return redirect('regex:search_form')

    directory = request.POST.get('directory')
    selected_pattern = request.POST.get('pattern')
    tools = request.POST.getlist('tools')

    errors = []
    results = []

    # Dizini doğrula
    if not directory or not os.path.isabs(directory) or not os.path.isdir(directory):
        errors.append(f"Geçersiz dizin: {directory}")

    # Deseni doğrula ve meta verileri al
    pattern_info = sensitive_patterns.get(selected_pattern)
    if not pattern_info:
        errors.append(f"Desen bulunamadı: {selected_pattern}")
    else:
        regex_str = pattern_info['pattern']
        try:
            compiled = re.compile(regex_str)
        except re.error as e:
            errors.append(f"Geçersiz regex deseni: {e}")

    if errors:
        return render(request, 'regex/search_results.html', {'errors': errors})

    # Python tabanlı tarama
    for root, _, files in os.walk(directory):
        for fname in files:
            fpath = os.path.join(root, fname)
            try:
                with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
                    for i, line in enumerate(f, start=1):
                        if compiled.search(line):
                            results.append({
                                'file': fpath,
                                'line': i,
                                'text': line.strip(),
                                'pattern': selected_pattern,
                                'category': pattern_info['category'],
                                'subcategory': pattern_info['subcategory'],
                                'policy': pattern_info['policy'],
                                'tool': 'python'
                            })
            except Exception as ex:
                results.append({'file': fpath, 'error': str(ex)})

    # Harici araçlarla tarama
    tool_map = {
        'grep': scan_with_grep,
        'ripgrep': scan_with_ripgrep,
        'ag': scan_with_ag,
        'ack': scan_with_ack,
    }
    for t in tools:
        func = tool_map.get(t)
        if not func:
            continue
        try:
            hits = func(directory, regex_str)
            for hit in hits:
                parts = hit.split(':', 2)
                if len(parts) == 3:
                    fpath, lineno, text = parts
                    results.append({
                        'file': fpath,
                        'line': int(lineno),
                        'text': text.strip(),
                        'pattern': selected_pattern,
                        'category': pattern_info['category'],
                        'subcategory': pattern_info['subcategory'],
                        'policy': pattern_info['policy'],
                        'tool': t
                    })
                else:
                    results.append({'match': hit, 'tool': t})
        except Exception as ex:
            results.append({'tool': t, 'error': str(ex)})

    # Sonuçları JSON dosyasına kaydet
    results_file = os.path.join(BASE_DIR, 'regex_results.json')
    try:
        with open(results_file, 'w', encoding='utf-8') as rf:
            json.dump(results, rf, ensure_ascii=False, indent=2)
    except Exception as ex:
        errors.append(f"Sonuç kaydedilemedi: {ex}")

    context = {
        'results': results,
        'errors': errors,
        'selected_tools': tools,
        'pattern': selected_pattern,
    }
    return render(request, 'regex/search_results.html', context)

# Maskeleme ve şifreleme fonksiyonları
def generate_encryption_key(password: str, salt: bytes = None):
    """Şifreleme anahtarı oluşturur."""
    if salt is None:
        salt = os.urandom(16)
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,
        salt=salt,
        iterations=100000,
    )
    key = base64.urlsafe_b64encode(kdf.derive(password.encode()))
    return key, salt

def encrypt_text(text: str, password: str):
    """Metni şifreler."""
    key, salt = generate_encryption_key(password)
    f = Fernet(key)
    encrypted_data = f.encrypt(text.encode())
    return base64.urlsafe_b64encode(salt + encrypted_data).decode()

def decrypt_text(encrypted_text: str, password: str):
    """Şifrelenmiş metni çözer."""
    try:
        data = base64.urlsafe_b64decode(encrypted_text.encode())
        salt = data[:16]
        encrypted_data = data[16:]
        key, _ = generate_encryption_key(password, salt)
        f = Fernet(key)
        decrypted_data = f.decrypt(encrypted_data)
        return decrypted_data.decode()
    except Exception as e:
        logger.error(f"Decryption error: {str(e)}")
        return None

def mask_sensitive_data(text: str, mask_type: str = 'asterisk'):
    """Hassas veriyi maskeler."""
    patterns_to_mask = {
        'tc_kimlik': r'\b[1-9][0-9]{10}\b',
        'email': r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
        'phone': r'\b(?:\+90|0)?\s*?\(?5\d{2}\)?[\s-]?\d{3}[\s-]?\d{2}[\s-]?\d{2}',
        'credit_card': r'\b(?:\d[ -]*?){13,16}\b',
        'iban': r'TR\d{2}\s?(\d{4}[ ]?){5}\d{2}',
    }
    
    masked_text = text
    for pattern_name, pattern in patterns_to_mask.items():
        if mask_type == 'asterisk':
            masked_text = re.sub(pattern, lambda m: '*' * len(m.group()), masked_text)
        elif mask_type == 'hash':
            masked_text = re.sub(pattern, lambda m: hashlib.md5(m.group().encode()).hexdigest()[:8], masked_text)
        elif mask_type == 'partial':
            def partial_mask(match):
                value = match.group()
                if len(value) > 4:
                    return value[:2] + '*' * (len(value) - 4) + value[-2:]
                return '*' * len(value)
            masked_text = re.sub(pattern, partial_mask, masked_text)
    
    return masked_text

def get_enhanced_file_content(file_path, include_metadata=True):
    """Gelişmiş dosya içeriği çıkarma."""
    try:
        file_type = get_file_type(file_path)
        content_data = {
            'raw_content': '',
            'metadata': {},
            'pages': [],
            'extraction_method': 'basic'
        }
        
        if file_type == 'pdf':
            try:
                import PyPDF2
                import pdfplumber
                
                # PyPDF2 ile temel çıkarım
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    
                    if include_metadata:
                        content_data['metadata'] = {
                            'page_count': len(pdf_reader.pages),
                            'title': pdf_reader.metadata.get('/Title', 'N/A') if pdf_reader.metadata else 'N/A',
                            'author': pdf_reader.metadata.get('/Author', 'N/A') if pdf_reader.metadata else 'N/A',
                            'creator': pdf_reader.metadata.get('/Creator', 'N/A') if pdf_reader.metadata else 'N/A',
                            'creation_date': str(pdf_reader.metadata.get('/CreationDate', 'N/A')) if pdf_reader.metadata else 'N/A'
                        }
                    
                    # pdfplumber ile gelişmiş çıkarım
                    try:
                        with pdfplumber.open(file_path) as pdf:
                            for i, page in enumerate(pdf.pages):
                                page_content = page.extract_text()
                                tables = page.extract_tables()
                                
                                page_data = {
                                    'page_number': i + 1,
                                    'text': page_content or '',
                                    'tables': tables or [],
                                    'char_count': len(page_content) if page_content else 0
                                }
                                content_data['pages'].append(page_data)
                                content_data['raw_content'] += (page_content or '') + '\n'
                        
                        content_data['extraction_method'] = 'pdfplumber'
                    except ImportError:
                        # pdfplumber yoksa PyPDF2 kullan
                        for page in pdf_reader.pages:
                            page_text = page.extract_text()
                            content_data['raw_content'] += page_text + '\n'
                        content_data['extraction_method'] = 'pypdf2'
                        
            except Exception as e:
                logger.error(f"Enhanced PDF reading error {file_path}: {str(e)}")
                return get_file_content(file_path)  # Fallback to basic method
                
        elif file_type in ['docx', 'doc']:
            try:
                import docx
                doc = docx.Document(file_path)
                
                if include_metadata:
                    core_props = doc.core_properties
                    content_data['metadata'] = {
                        'title': core_props.title or 'N/A',
                        'author': core_props.author or 'N/A',
                        'created': str(core_props.created) if core_props.created else 'N/A',
                        'modified': str(core_props.modified) if core_props.modified else 'N/A',
                        'paragraph_count': len(doc.paragraphs),
                        'table_count': len(doc.tables)
                    }
                
                # Paragrafları çıkar
                for para in doc.paragraphs:
                    content_data['raw_content'] += para.text + '\n'
                
                # Tabloları çıkar
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    content_data['raw_content'] += '\n' + str(table_data) + '\n'
                    
                content_data['extraction_method'] = 'python-docx'
                
            except Exception as e:
                logger.error(f"Enhanced DOCX reading error {file_path}: {str(e)}")
                return get_file_content(file_path)
                
        elif file_type in ['xlsx', 'xls']:
            try:
                import pandas as pd
                
                # Tüm sheet'leri oku
                excel_file = pd.ExcelFile(file_path)
                
                if include_metadata:
                    content_data['metadata'] = {
                        'sheet_names': excel_file.sheet_names,
                        'sheet_count': len(excel_file.sheet_names)
                    }
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    content_data['raw_content'] += f'\n=== Sheet: {sheet_name} ===\n'
                    content_data['raw_content'] += df.to_string() + '\n'
                
                content_data['extraction_method'] = 'pandas'
                
            except Exception as e:
                logger.error(f"Enhanced Excel reading error {file_path}: {str(e)}")
                return get_file_content(file_path)
                
        elif file_type in ['pptx', 'ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                
                if include_metadata:
                    content_data['metadata'] = {
                        'slide_count': len(prs.slides),
                        'slide_layouts': len(prs.slide_layouts)
                    }
                
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    
                    content_data['raw_content'] += f'\n=== Slide {i+1} ===\n'
                    content_data['raw_content'] += '\n'.join(slide_text) + '\n'
                
                content_data['extraction_method'] = 'python-pptx'
                
            except Exception as e:
                logger.error(f"Enhanced PPTX reading error {file_path}: {str(e)}")
                return get_file_content(file_path)
        
        else:
            # Diğer dosya türleri için temel metod kullan
            content_data['raw_content'] = get_file_content(file_path) or ''
            content_data['extraction_method'] = 'basic'
        
        return content_data
        
    except Exception as e:
        logger.error(f"Enhanced file content extraction error {file_path}: {str(e)}")
        return {'raw_content': get_file_content(file_path) or '', 'metadata': {}, 'extraction_method': 'fallback'}
